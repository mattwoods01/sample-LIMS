from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from num2words import num2words
import pandas as pd
import numpy as np
from datetime import datetime

# Function to get a placeholder by type
def get_placeholder_by_type(slide, placeholder_type):
    """
    get placeholder type.  different types of placeholders exist for powerpoints.  Some can be text, images, text boxes, etc.
    """
    for placeholder in slide.placeholders:
        if placeholder.placeholder_format.idx == placeholder_type:
            return placeholder
    raise KeyError(f"No placeholder with type {placeholder_type}")

def pluritest_generator(template_path, form_df, output_path, input_df, image_directory):
    """
    pluritest generator function.  sets up format for pluritest powerpoint.
    """

    # Load the existing PowerPoint template
    prs = Presentation(template_path)
    
    print(output_path)
   
    # Update the first slide with form data
    slide = prs.slides[0]

    # Assuming the first slide layout has the placeholders we need
    client_name = form_df['Customer Affiliation'].values[0]
    quote_no = form_df['Project #'].values[0]
    date = datetime.now().strftime("%d %B %Y")
    prepared_by = form_df['Technical Project Lead'].values[0]

    text_to_add = (
        f"Client Name: {client_name}\n"
        f"Quote No: {quote_no}\n"
        f"Date: {date}\n"
        f"Prepared by: {prepared_by}"
    )
    
    body_placeholder_type = 13   # Typically body placeholder is of type 1, update if necessary

    try:
        placeholder = slide.placeholders[body_placeholder_type]
        placeholder.text_frame.text = text_to_add

    except KeyError as e:
        print(f"Placeholder with index {body_placeholder_type} not found: {e}")

    slide = prs.slides[1]

    num_samples = form_df['# Samples'].values[0]
    num_samples_word = num2words(num_samples)

    project_summary_text = (
        f"{client_name} (Client) is interested in Services provided by the Life Technologies Corporation in the analysis of {num_samples_word} ({num_samples}) Client provided sample(s) using the PluriTest™ Service. "
        f"In this assay, 36,000 transcripts and variants against a >450 sample reference set are assessed for gene expression analysis. "
    )

    project_summary_placeholder_idx = 11  # Update this index based on your slide layout
    try:
        placeholder = slide.placeholders[project_summary_placeholder_idx]
        placeholder.text_frame.text = project_summary_text

    except KeyError as e:
        print(f"Placeholder with index {project_summary_placeholder_idx} not found: {e}")

    # Add table to the third slide using a placeholder
    slide = prs.slides[2]  # Select the third slide (0-based index)
    filtered_df = input_df[input_df['Project'] != '']

    print(filtered_df)

    table_df = input_df[['#', 'Sample ID', 'Pluripotency Pass/Fail', 'PluriCor', 'NovelCor']]
    table_df.replace(['', 'nan'], np.nan, inplace=True)
    table_df = table_df.dropna(subset=['#'])

    table_df = table_df.rename(columns={
        '#': 'Sample',
        'Sample ID': 'Sample ID',
        'Pluripotency Pass/Fail': 'PluriTest Result',
        'PluriCor': 'PluriCor',
        'NovelCor': 'NovelCor'
    })

    if (filtered_df['Pluripotency Pass/Fail'] == 'Pass').all():
        result = "All sample(s) were found to be pluripotent."

    else:
        passed_samples = filtered_df.loc[filtered_df['Pluripotency Pass/Fail'] == 'Pass', '#'].tolist()
        result = f"Sample(s) {', '.join(passed_samples)} were found to be pluripotent."

    project_summary_text = (
        f"{result}\n"
        "The raw gene expression data for these sample(s) will be delivered to the client."
    )

    project_summary_placeholder_idx = 14  # Update this index based on your slide layout
    try:
        placeholder = slide.placeholders[project_summary_placeholder_idx]
        placeholder.text_frame.text = project_summary_text

    except KeyError as e:
        print(f"Placeholder with index {project_summary_placeholder_idx} not found: {e}")

    placeholder_idx = 17  # Update this index based on your slide layout
    try:
        placeholder = slide.placeholders[placeholder_idx]

        # Clear the placeholder to ensure it's empty
        placeholder.text = ''

        # Calculate rows and columns
        rows, cols = table_df.shape
        rows += 1  # Add one for the header

        # Add table to the placeholder
        table = placeholder.insert_table(rows, cols).table

        # Set uniform column widths
        table.columns[0].width = Inches(1.3)
        table.columns[1].width = Inches(2)
        table.columns[2].width = Inches(0.9)
        table.columns[3].width = Inches(0.9)
        table.columns[4].width = Inches(0.9)

        # Set row height for header and all rows
        row_height = Inches(0.3)
        for i in range(rows):
            table.rows[i].height = row_height

        # Add header row with custom formatting
        for col, column_name in enumerate(table_df.columns):
            cell = table.cell(0, col)
            cell.text = column_name
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(79, 79, 79)  # Dark gray background
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.bold = True
                    run.font.size = Pt(10)
                    run.font.color.rgb = RGBColor(255, 255, 255)  # White text
                    run.font.name = 'Calibri'
                paragraph.alignment = PP_ALIGN.CENTER  # Center-align the header text

        # Add data rows with alternating colors
        light_grey = RGBColor(192, 192, 192)
        white = RGBColor(255, 255, 255)

        for row in range(1, rows):
            row_color = white if row % 2 == 1 else light_grey  # Alternate colors
            for col in range(cols):
                cell = table.cell(row, col)
                cell.text = str(table_df.iat[row-1, col])
                cell.fill.solid()
                cell.fill.fore_color.rgb = row_color  # Set the background color for the row

                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Pt(9)
                        run.font.color.rgb = RGBColor(0, 0, 0)  # Black text
                    paragraph.alignment = PP_ALIGN.CENTER  # Center-align the cell text

    except KeyError as e:
        print(f"Placeholder with index {placeholder_idx} not found: {e}")

    # Add an image to a placeholder on the same slide
    slide = prs.slides[2]
    placeholder_idx = 16  # Update this index based on your slide layout
    try:
        placeholder = slide.placeholders[placeholder_idx]
        placeholder.insert_picture(image_directory)

    except KeyError as e:
        print(f"Placeholder with index {placeholder_idx} not found: {e}")

    # Save the presentation
    prs.save(output_path)
    print(f"Presentation saved to {output_path}")

